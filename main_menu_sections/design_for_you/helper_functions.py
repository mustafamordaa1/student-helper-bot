import asyncio
import io
import datetime
from typing import Tuple
import uuid
import logging
import os
import openpyxl
from PIL import Image
from pptx import Presentation
import aiohttp
from config import DESIGNS_FOR_FEMALE_FILE, DESIGNS_FOR_MALE_FILE
from utils.database import execute_query, get_data
from utils.user_management import get_user_data
import tempfile
from pdf2image import convert_from_path
import os
import subprocess

logger = logging.getLogger(__name__)


def load_design_options(gender: str):
    try:
        if gender == "male":
            FILE_GENDER = DESIGNS_FOR_MALE_FILE
        else:
            FILE_GENDER = DESIGNS_FOR_FEMALE_FILE
        workbook = openpyxl.load_workbook(FILE_GENDER)
        sheet = workbook.active
        return [
            (row[0], row[1]) for row in sheet.iter_rows(min_row=2, values_only=True)
        ]
    except Exception as e:
        logger.error(f"Error loading design options: {e}")
        return []


def process_powerpoint_design(ppt_file, user_name):
    try:
        prs = Presentation(ppt_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text.find("Your Name") != -1:
                        shape.text = shape.text.replace("Your Name", user_name)
        modified_ppt = f"temp_{user_name}.pptx"
        prs.save(modified_ppt)

        image_path = f"temp_{user_name}.png"
        convert_ppt_to_image(modified_ppt, image_path)

        os.remove(modified_ppt)
        return image_path
    except Exception as e:
        logger.error(f"Error processing PowerPoint design: {e}")
        raise

def convert_ppt_to_image(ppt_file_path, img_path, dpi=300, image_format="png"):
    """
    Converts a PPTX file to images by first converting it to a PDF, then converting each PDF page to an image.
    
    Parameters:
    - ppt_file_path: Path to the input PPTX file
    - output_dir: Path to the directory where images will be saved
    - dpi: DPI setting for converting PDF to images (default is 300)
    - image_format: Format to save images (default is PNG)
    """
    
    # Step 1: Convert PPTX to PDF using LibreOffice
    with tempfile.TemporaryDirectory() as temp_dir:
        pdf_filename = os.path.splitext(os.path.basename(ppt_file_path))[0] + ".pdf"
        pdf_path = os.path.join(temp_dir, pdf_filename)
        
        # Run the LibreOffice command to convert PPTX to PDF
        result = subprocess.run([
            "libreoffice", "--headless", "--invisible", "--convert-to", "pdf",
            "--outdir", temp_dir, ppt_file_path
        ], check=True)
        
        # Check if the PDF was created successfully
        if not os.path.exists(pdf_path):
            raise FileNotFoundError("PDF conversion failed.")
        print(f"Converted PPTX to PDF: {pdf_path}")
        
        # Step 2: Convert PDF to images
        images = convert_from_path(pdf_path, dpi=dpi, fmt=image_format)
        
        # Step 3: Save images to the specified output directory
        for i, img in enumerate(images):
            img.save(img_path, image_format.upper())

            break
        
        print(f"PDF converted to images. Images saved in {img_path}")
        
        return img_path

async def download_image(image_url):
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(image_url) as response:
                response.raise_for_status()
                file_name = f"{uuid.uuid4()}.png"
                file_path = os.path.join("temp", file_name)
                os.makedirs(
                    "temp", exist_ok=True
                )  # Create temp directory if it doesn't exist
                with open(file_path, "wb") as f:
                    async for chunk in response.content.iter_chunked(1024):
                        f.write(chunk)
                return file_path
    except aiohttp.ClientError as e:
        logger.error(f"Error downloading image: {e}")
        raise


async def check_user_ai_limit(user_id: int) -> Tuple[bool, int, int]:
    """Checks if user has reached their daily AI design limit, considering subscription type."""
    today = datetime.date.today()
    today_start = datetime.datetime(today.year, today.month, today.day, 0, 0, 0)

    query = """
        SELECT COUNT(*) 
        FROM ai_image_usage 
        WHERE user_id = ? 
        AND usage_time >= ?
    """
    params = (user_id, today_start)
    usage_count = get_data(query, params)[0][0]

    user_data = await asyncio.to_thread(get_user_data, user_id)
    subscription_type = user_data.get("type_of_last_subscription")

    daily_limit = get_daily_ai_limit(subscription_type)

    is_not_subscribed = daily_limit == 1
    images_left = daily_limit - usage_count
    is_allowed = images_left > 0

    return is_allowed, usage_count, images_left, is_not_subscribed


def get_daily_ai_limit(subscription_type: str) -> int:
    """Returns the daily AI image limit based on the subscription type."""
    if (
        subscription_type == "تجربة مجانية الساعية"
        or subscription_type == "تجربة مجانية"
    ):
        return 1
    else:  # Subscripted
        return 2


async def update_user_ai_usage(user_id: int):
    query = """
        INSERT INTO ai_image_usage (user_id, usage_time) 
        VALUES (?, ?)
    """
    params = (user_id, datetime.datetime.now())
    await asyncio.to_thread(execute_query, query, params)  # Make this awaitable
